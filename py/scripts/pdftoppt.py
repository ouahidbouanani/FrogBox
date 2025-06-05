import io
from pdf2image import convert_from_bytes
from pptx import Presentation
from pptx.util import Inches
import os

# Ajouter le chemin de Poppler
poppler_path = r'.\poppler-24.08.0\Library\bin'
os.environ["PATH"] += os.pathsep + poppler_path

def convert_pdf_to_pptx(pdf_bytes):
    try:
        images = convert_from_bytes(pdf_bytes)
    except Exception as e:
        return {"error": str(e)}

    prs = Presentation()

    for img in images:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Layout vide
        left = top = Inches(0)
        width = Inches(10)
        height = Inches(7.5)
        
        img_io = io.BytesIO()
        img.save(img_io, format='PNG')
        img_io.seek(0)

        slide.shapes.add_picture(img_io, left, top, width, height)

    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)

    return ppt_io

